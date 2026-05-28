#!/usr/bin/env python3
"""
extract_pptx.py — Phase 4 extractor for PPTX files. Critically, this
preserves speaker notes — they often contain more semantic content than the
slides themselves, and most off-the-shelf converters drop them. Equally
critical, this extracts embedded pictures to disk via save_image_safe() —
without that, diagram-only slides ("Пример 1", "Пример 2", …) collapse to
just the header and the second-session agent has no way to know what the
slide was about.

CLI:
    extract_pptx.py <input_pptx> <output_md>
                    [--doc-id doc-NNN]
                    [--source-rel relative/path/in/corpus.pptx]
                    [--assets-dir <abs_path>]
                    [--assets-rel <rel_prefix_from_md>]
                    [--no-extract-images]

Output layout (per slide):

    ## Slide N: <best-effort title>

    <body text from text frames, in shape order>

    ![Slide N picture M](../assets/doc-NNN-slideNN-imgM.png)

    | header | cells |
    | ---    | ---   |
    | ...    | ...   |

    ### Notes

    <speaker notes>

    ---  (slide separator)

The body intentionally interleaves shape contents in document order rather
than re-ordering by position — readers care about semantics, not layout.

Image handling:
    - By default, embedded PICTURE shapes (and pictures nested in GROUP
      shapes one level deep) are extracted via save_image_safe() to the
      kb's `assets/` directory and referenced from the slide body as
      Markdown image links. Animations, palette/RGBA-as-JPEG mismatches,
      and tiny decorative pixels are filtered automatically.
    - WMF/EMF (Visio paste, "vector clipart") cannot be rendered by the
      vision encoder; they are skipped and recorded in `warnings`.
    - `--no-extract-images` reverts to the old count-only behaviour.
"""
from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

from _common import (  # noqa: E402
    add_assets_args,
    clean_whitespace,
    count_tokens,
    emit_failure,
    emit_success,
    resolve_assets_dir,
    sanitize_heading,
    save_image_safe,
    sha256_of,
    today_iso,
    tool_version_string,
    validate_source_rel,
    write_md,
)


EXTRACTOR_NAME = "python-pptx"


def _try_import():
    try:
        from pptx import Presentation  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
        return Presentation, MSO_SHAPE_TYPE
    except Exception as e:
        emit_failure(f"python-pptx unavailable: {e}")
        sys.exit(1)


def _is_picture_like(shape) -> bool:
    """True if `shape` carries an embedded raster the way a Picture shape
    does, even when its `shape_type` is PLACEHOLDER rather than PICTURE.
    `PlaceholderPicture` inherits from `Picture` in python-pptx — both
    expose `shape.image.blob`, but the placeholder is classified by
    placeholder type and would otherwise fall through the PICTURE branch
    silently. Catches the common "title + image" slide layout that
    BPMN/UML tutorial decks use to show one diagram per slide."""
    if shape is None:
        return False
    try:
        from pptx.shapes.picture import Picture  # type: ignore
    except Exception:
        Picture = None  # noqa: N806
    if Picture is not None and isinstance(shape, Picture):
        return True
    # Fallback duck-typing — some pptx variants expose `.image` on objects
    # that don't subclass Picture.
    img = getattr(shape, "image", None)
    return img is not None and hasattr(img, "blob")


# Placeholder types that contain layout chrome (slide number, header, footer,
# date). These are noise for LLM consumption — skip them.
# Values from pptx.enum.text.PP_PLACEHOLDER:
#   13 SLIDE_NUMBER, 14 HEADER, 15 FOOTER, 16 DATE
_SKIP_PLACEHOLDER_TYPES = {13, 14, 15, 16}


def _is_chrome_placeholder(shape) -> bool:
    try:
        if not shape.is_placeholder:
            return False
        return int(shape.placeholder_format.type) in _SKIP_PLACEHOLDER_TYPES
    except Exception:
        return False


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    try:
        tf = shape.text_frame
    except Exception:
        return ""
    parts = []
    for para in tf.paragraphs:
        line = "".join(run.text for run in para.runs)
        if not line:
            # Sometimes runs are empty but para.text has content (auto-numbered, etc.)
            line = para.text or ""
        parts.append(line)
    return "\n".join(parts).rstrip()


def _table_md(shape) -> str:
    """Renders a pptx table as a Markdown pipe-table. Cells are joined with
    a space when they contain multiple paragraphs to avoid breaking the
    pipe-table format."""
    table = shape.table
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            tf = cell.text_frame
            txt = " ".join(p.text for p in tf.paragraphs if p.text).replace("\n", " ").strip()
            # Markdown pipes inside cells break the table; escape them.
            txt = txt.replace("|", "\\|")
            cells.append(txt)
        rows.append(cells)
    if not rows:
        return ""
    ncols = max(len(r) for r in rows)
    rows = [r + [""] * (ncols - len(r)) for r in rows]
    sep = ["---"] * ncols
    lines = ["| " + " | ".join(rows[0]) + " |",
             "| " + " | ".join(sep) + " |"]
    for r in rows[1:]:
        lines.append("| " + " | ".join(r) + " |")
    return "\n".join(lines)


def _best_title(slide) -> str | None:
    # 1) explicit title placeholder
    for shape in slide.shapes:
        try:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                t = _shape_text(shape).strip().split("\n", 1)[0]
                if t:
                    return t[:200]
        except Exception:
            pass
    # 2) first text shape's first line, if it looks like a heading
    for shape in slide.shapes:
        try:
            if shape.has_text_frame:
                t = _shape_text(shape).strip().split("\n", 1)[0]
                if t and len(t) < 200:
                    return t
        except Exception:
            pass
    return None


def _picture_caption(shape) -> str:
    """Pull a human-friendly caption from a PICTURE shape. Tries alt_text
    first (set when authors describe the image for accessibility), falls
    back to the shape name (often "Picture 12" — uninformative but stable),
    then to "" so the caller can decide what to write."""
    for attr in ("alt_text", "name"):
        try:
            val = getattr(shape, attr, "") or ""
            val = val.strip()
            if val and val.lower() not in ("picture", "image"):
                return sanitize_heading(val, maxlen=120)
        except Exception:
            pass
    return ""


def _save_picture_shape(
    shape,
    doc_id: str,
    slide_idx: int,
    pic_idx: int,
    assets_dir: Path,
    assets_rel_prefix: str,
    dedup_cache: dict,
    warnings: list[str],
) -> tuple[str, str] | None:
    """Extract a single PICTURE shape's bytes, run them through
    save_image_safe(), and return (markdown_image_ref, asset_relpath) on
    success. Returns None when the image was filtered (too small, unviewable
    format) — caller decides whether to emit a fallback marker.

    The filename pattern is `<doc_id>-slideNN-imgM.<ext>`. Deduplication is
    intra-document only: a logo repeated on every slide saves disk by
    reusing the first occurrence, but two different documents that share a
    bitwise-identical image still keep separate copies (cross-doc dedup
    would require a manifest-level pass).
    """
    try:
        image = shape.image  # type: ignore[attr-defined]
        blob = image.blob
        ext = (image.ext or "png").lower()
    except Exception as e:
        warnings.append(
            f"slide {slide_idx} picture {pic_idx}: cannot read image blob ({e})"
        )
        return None

    filename = f"{doc_id}-slide{slide_idx:02d}-img{pic_idx}.{ext}"
    target = assets_dir / filename
    result = save_image_safe(
        blob, target, dedup_cache=dedup_cache, source_ext=ext,
    )
    if not result:
        if result.reason == "skipped_format":
            warnings.append(
                f"slide {slide_idx} picture {pic_idx}: skipped non-viewable "
                f"format ({ext}) — vector/metafile images cannot be rendered "
                "by the vision encoder"
            )
        elif result.reason == "skipped_corrupt":
            warnings.append(
                f"slide {slide_idx} picture {pic_idx}: corrupt or unreadable "
                "image bytes — skipped"
            )
        # too-small / duplicate-with-no-prior-path: silent, by design.
        return None

    rel = f"{assets_rel_prefix}/{result.path.name}"
    caption = _picture_caption(shape) or f"Slide {slide_idx} picture {pic_idx}"
    md_ref = f"![{caption}]({rel})"
    return md_ref, rel


def extract(
    input_pptx: Path,
    doc_id: str = "doc-000",
    assets_dir: Path | None = None,
    assets_rel_prefix: str = "../assets",
    extract_images: bool = True,
) -> tuple[str, dict, list[str]]:
    Presentation, MSO_SHAPE_TYPE = _try_import()
    warnings: list[str] = []
    extras: dict = {
        "slides": 0,
        "has_notes": False,
        "notes_chars": 0,
        "inline_images": 0,
        "has_tables": False,
        "has_charts": False,
    }

    try:
        prs = Presentation(str(input_pptx))
    except Exception as e:
        raise RuntimeError(f"failed to open pptx: {e}") from e

    extras["slides"] = len(prs.slides)
    pieces: list[str] = []
    n_images = 0
    n_tables = 0
    n_charts = 0
    total_notes = 0
    asset_relpaths: list[str] = []
    dedup_cache: dict = {}

    do_extract = extract_images and assets_dir is not None
    if do_extract:
        try:
            assets_dir.mkdir(parents=True, exist_ok=True)
        except OSError as e:
            warnings.append(
                f"cannot create assets dir {assets_dir}: {e} — falling back to count-only"
            )
            do_extract = False

    for idx, slide in enumerate(prs.slides, start=1):
        title = _best_title(slide)
        if title:
            title = sanitize_heading(title)
        header = f"## Slide {idx}"
        if title:
            header += f": {title}"
        slide_parts: list[str] = [header, ""]

        title_shape_id = None
        # Find the title shape id so we don't print it twice.
        for shape in slide.shapes:
            try:
                if shape.is_placeholder and shape.placeholder_format.idx == 0:
                    title_shape_id = shape.shape_id
                    break
            except Exception:
                pass

        body_text_pieces: list[str] = []
        table_pieces: list[str] = []
        slide_pic_idx = 0  # 1-based counter for image refs within this slide

        for shape in slide.shapes:
            try:
                stype = shape.shape_type
            except Exception:
                stype = None

            if _is_chrome_placeholder(shape):
                continue

            # Pictures (regular + placeholder-pictures, e.g. "Picture 2"
            # filling a content placeholder on diagram-only slides).
            if stype == MSO_SHAPE_TYPE.PICTURE or _is_picture_like(shape):
                n_images += 1
                slide_pic_idx += 1
                if do_extract:
                    saved = _save_picture_shape(
                        shape, doc_id, idx, slide_pic_idx,
                        assets_dir, assets_rel_prefix, dedup_cache, warnings,
                    )
                    if saved is not None:
                        md_ref, relpath = saved
                        body_text_pieces.append(md_ref)
                        if relpath not in asset_relpaths:
                            asset_relpaths.append(relpath)
                    else:
                        body_text_pieces.append(
                            f"*[Slide {idx} picture {slide_pic_idx} skipped]*"
                        )
                continue

            # Tables.
            if getattr(shape, "has_table", False):
                n_tables += 1
                md = _table_md(shape)
                if md:
                    table_pieces.append(md)
                continue

            # Charts — we can't render them as markdown; record presence.
            if stype == MSO_SHAPE_TYPE.CHART or getattr(shape, "has_chart", False):
                n_charts += 1
                body_text_pieces.append("*(chart)*")
                continue

            # Text frames.
            if getattr(shape, "has_text_frame", False):
                if title_shape_id is not None and shape.shape_id == title_shape_id:
                    # Already used as slide header; skip duplicate.
                    continue
                txt = _shape_text(shape)
                if txt.strip():
                    body_text_pieces.append(txt.strip())
                continue

            # Grouped shapes — recurse one level to catch text inside groups.
            if stype == MSO_SHAPE_TYPE.GROUP:
                try:
                    for gsub in shape.shapes:
                        if getattr(gsub, "has_text_frame", False):
                            t = _shape_text(gsub).strip()
                            if t:
                                body_text_pieces.append(t)
                        elif (gsub.shape_type == MSO_SHAPE_TYPE.PICTURE
                              or _is_picture_like(gsub)):
                            n_images += 1
                            slide_pic_idx += 1
                            if do_extract:
                                saved = _save_picture_shape(
                                    gsub, doc_id, idx, slide_pic_idx,
                                    assets_dir, assets_rel_prefix,
                                    dedup_cache, warnings,
                                )
                                if saved is not None:
                                    md_ref, relpath = saved
                                    body_text_pieces.append(md_ref)
                                    if relpath not in asset_relpaths:
                                        asset_relpaths.append(relpath)
                                else:
                                    body_text_pieces.append(
                                        f"*[Slide {idx} picture "
                                        f"{slide_pic_idx} skipped]*"
                                    )
                except Exception:
                    pass

        if body_text_pieces:
            slide_parts.append("\n\n".join(body_text_pieces))
        for tbl in table_pieces:
            slide_parts.append("")
            slide_parts.append(tbl)

        # Speaker notes.
        notes_text = ""
        if slide.has_notes_slide:
            try:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
            except Exception:
                notes_text = ""
        if notes_text:
            extras["has_notes"] = True
            total_notes += len(notes_text)
            slide_parts.append("")
            slide_parts.append("### Notes")
            slide_parts.append("")
            slide_parts.append(notes_text)

        pieces.append("\n".join(slide_parts).rstrip())

    body = "\n\n---\n\n".join(pieces).strip() + "\n"
    body = clean_whitespace(body)

    extras["notes_chars"] = total_notes
    extras["inline_images"] = n_images
    extras["has_tables"] = n_tables > 0
    extras["has_charts"] = n_charts > 0
    if asset_relpaths:
        extras["assets"] = asset_relpaths
    if n_charts:
        warnings.append(f"{n_charts} chart(s) skipped (rendered as *(chart)* placeholder)")

    # Collect slide titles only, stripping the "Slide N: " chrome prefix so
    # `headings` is consistent with how other extractors record headings
    # (the actual title text, not the structural anchor).
    headings = []
    for line in body.split("\n"):
        if line.startswith("## ") and len(line) < 200:
            heading = line[3:]
            # "Slide N: Title" → "Title"; bare "Slide N" → skip (no title).
            m = re.match(r"^Slide \d+(?::\s*(.+))?$", heading)
            if m:
                title = (m.group(1) or "").strip()
                if not title:
                    continue
                heading = title
            sanitized = sanitize_heading(heading)
            if sanitized:
                headings.append(sanitized)
            if len(headings) >= 15:
                break
    extras["headings"] = headings

    return body, extras, warnings


def main() -> int:
    ap = argparse.ArgumentParser(description="Extract PPTX → Markdown via python-pptx.")
    ap.add_argument("input_pptx")
    ap.add_argument("output_md")
    ap.add_argument("--doc-id", default="doc-000")
    ap.add_argument("--source-rel", default=None)
    add_assets_args(ap)
    args = ap.parse_args()

    in_path = Path(args.input_pptx).expanduser().resolve()
    out_path = Path(args.output_md).expanduser().resolve()
    source_rel = args.source_rel or in_path.name
    try:
        source_rel = validate_source_rel(source_rel)
    except ValueError as e:
        emit_failure(f"invalid --source-rel: {e}")
        return 1

    if not in_path.is_file():
        emit_failure(f"input not found: {in_path}")
        return 1

    # Default layout: <kb_dir>/docs/*.md + <kb_dir>/assets/.
    assets_dir: Path | None = resolve_assets_dir(args, out_path)

    try:
        body, extras, warnings = extract(
            in_path,
            doc_id=args.doc_id,
            assets_dir=assets_dir,
            assets_rel_prefix=args.assets_rel,
            extract_images=not args.no_extract_images,
        )
    except Exception as e:
        emit_failure(f"extraction failed: {e}", extra={"input": str(in_path)})
        return 1

    fm = {
        "id": args.doc_id,
        "source": source_rel,
        "source_type": "pptx",
        "source_sha256": sha256_of(in_path),
        "extraction_method": f"{EXTRACTOR_NAME}@{tool_version_string()}",
        "extraction_date": today_iso(),
        "slides": extras.get("slides", 0),
        "has_notes": extras.get("has_notes", False),
        "notes_chars": extras.get("notes_chars", 0),
        "inline_images": extras.get("inline_images", 0),
        "has_tables": extras.get("has_tables", False),
        "has_charts": extras.get("has_charts", False),
        "headings": extras.get("headings", []),
        "tokens_estimated": count_tokens(body),
        "warnings": warnings,
    }
    if extras.get("assets"):
        fm["assets"] = extras["assets"]
    write_md(out_path, fm, body)
    emit_success(out_path, body, warnings, extra={
        "slides": extras.get("slides", 0),
        "has_notes": extras.get("has_notes", False),
        "assets_extracted": len(extras.get("assets", [])),
    })
    return 0


if __name__ == "__main__":
    sys.exit(main())
